# analysis/text_extraction.py
"""
Text extraction from various file formats.
"""

import os
from pathlib import Path
from typing import List, Dict, Any, Tuple, Optional
import re

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from PIL import Image
    import pytesseract
except ImportError:
    Image = None
    pytesseract = None


def extract_text_from_pdf(file_path: str) -> Tuple[str, Optional[str]]:
    """
    Extract text from a PDF file.
    Returns (extracted_text, error_message)
    """
    if PyPDF2 is None:
        return "", "PyPDF2 library not installed"
    
    text = ""
    try:
        with open(file_path, "rb") as f:
            pdf_reader = PyPDF2.PdfReader(f)
            
            if len(pdf_reader.pages) == 0:
                return "", "PDF has no pages"
            
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            
            if not text.strip():
                return "", "PDF appears to be empty or contains only images (no extractable text)"
                
    except Exception as e:
        error_msg = f"Error extracting PDF: {str(e)}"
        return "", error_msg
    
    return text, None


def extract_text_from_pptx(file_path: str) -> Tuple[str, Optional[str], Optional[List[Dict[str, Any]]]]:
    """
    Extract text from a PowerPoint file with structured slide data.
    Returns (extracted_text, error_message, structured_slides)
    - structured_slides: List of dicts with 'title', 'body', 'is_learning_objectives', 'is_key_ideas'
    """
    if Presentation is None:
        return "", "python-pptx library not installed", None
    
    text = ""
    structured_slides = []
    
    try:
        prs = Presentation(file_path)
        
        if len(prs.slides) == 0:
            return "", "PPTX has no slides", None
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_title = ""
            slide_body = []
            slide_text = ""
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    shape_text = shape.text.strip()
                    if shape_text:
                        slide_text += shape_text + "\n"
                        
                        # Try to identify title (usually first text box or title placeholder)
                        if not slide_title and len(shape_text) < 100:
                            # Check if it's likely a title (short, on first shape, or in title placeholder)
                            if shape.shape_type == 1 or (slide_idx == 0 and len(slide_body) == 0):  # Title placeholder
                                slide_title = shape_text
                            else:
                                slide_body.append(shape_text)
                        else:
                            slide_body.append(shape_text)
            
            # If no explicit title found, use first line
            if not slide_title and slide_body:
                first_line = slide_body[0].split('\n')[0]
                if len(first_line) < 100:
                    slide_title = first_line
                    slide_body = slide_body[1:] if len(slide_body) > 1 else []
            
            body_text = "\n".join(slide_body)
            
            # Detect learning objectives and key ideas slides
            title_lower = slide_title.lower()
            body_lower = body_text.lower()
            
            is_learning_objectives = any(keyword in title_lower or keyword in body_lower 
                                       for keyword in ["learning objectives", "by the end", "you will", 
                                                      "objectives", "goals", "outcomes"])
            
            is_key_ideas = any(keyword in title_lower or keyword in body_lower 
                             for keyword in ["key ideas", "summary", "takeaways", "main points", 
                                           "important points", "key concepts", "recap"])
            
            structured_slides.append({
                "title": slide_title,
                "body": body_text,
                "full_text": slide_text,
                "is_learning_objectives": is_learning_objectives,
                "is_key_ideas": is_key_ideas,
                "slide_index": slide_idx,
            })
            
            text += slide_text + "\n\n"
        
        if not text.strip():
            return "", "PPTX appears to be empty or contains only images (no extractable text)", None
            
    except Exception as e:
        error_msg = f"Error extracting PPTX: {str(e)}"
        return "", error_msg, None
    
    return text, None, structured_slides


def extract_text_from_image(file_path: str) -> Tuple[str, Optional[str]]:
    """
    Extract text from an image using OCR.
    Returns (extracted_text, error_message)
    """
    if Image is None or pytesseract is None:
        return "", "PIL/pytesseract libraries not installed"
    
    text = ""
    try:
        img = Image.open(file_path)
        text = pytesseract.image_to_string(img)
        
        if not text.strip():
            return "", "No text found in image (may be blank or unreadable)"
            
    except Exception as e:
        error_msg = f"Error extracting image: {str(e)}"
        return "", error_msg
    
    return text, None


def extract_all_text(file_paths: List[str]) -> Tuple[Dict[str, str], List[str], Dict[str, List[Dict[str, Any]]]]:
    """
    Extract text from all files with structured slide data.
    Returns (results_dict, error_messages, structured_data)
    - results_dict: mapping file_path -> extracted_text
    - error_messages: list of error messages for files that failed
    - structured_data: mapping file_path -> list of structured slides (for PPTX files)
    """
    results = {}
    errors = []
    structured_data = {}
    
    for file_path in file_paths:
        if not os.path.exists(file_path):
            errors.append(f"{os.path.basename(file_path)}: File not found")
            continue
        
        path = Path(file_path)
        ext = path.suffix.lower()
        
        text = ""
        error = None
        slides = None
        
        if ext == ".pdf":
            text, error = extract_text_from_pdf(file_path)
        elif ext in [".pptx", ".ppt"]:
            text, error, slides = extract_text_from_pptx(file_path)
            if slides:
                structured_data[file_path] = slides
        elif ext in [".png", ".jpg", ".jpeg"]:
            text, error = extract_text_from_image(file_path)
        else:
            error = f"Unsupported file type: {ext}"
        
        if error:
            errors.append(f"{os.path.basename(file_path)}: {error}")
        elif text:
            results[file_path] = text
        else:
            errors.append(f"{os.path.basename(file_path)}: No text extracted")
    
    return results, errors, structured_data

